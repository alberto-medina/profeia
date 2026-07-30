"""
Servicio de exportacion: genera PDF, PPTX y paquetes ZIP a partir del
contenido pedagogico.
"""

import json
from io import BytesIO
from pathlib import Path
from textwrap import wrap
from urllib.parse import urlparse
from uuid import UUID
from zipfile import ZIP_DEFLATED, ZipFile

import httpx
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.services.servicio_storage import guardar_recurso_docente


RUTA_FUENTES_WINDOWS = Path("C:/Windows/Fonts")
EXTENSIONES_IMAGEN_VALIDAS = {".png", ".jpg", ".jpeg", ".webp"}


def _limpiar_texto_pdf(texto: str) -> str:
    reemplazos = {
        "\\": "\\\\",
        "(": "\\(",
        ")": "\\)",
    }
    texto_limpio = str(texto)
    for original, reemplazo in reemplazos.items():
        texto_limpio = texto_limpio.replace(original, reemplazo)
    return texto_limpio


def _fuente_pdf(tamano: int, negrita: bool = False):
    nombres = (
        [RUTA_FUENTES_WINDOWS / "arialbd.ttf", RUTA_FUENTES_WINDOWS / "segoeuib.ttf"]
        if negrita
        else [RUTA_FUENTES_WINDOWS / "arial.ttf", RUTA_FUENTES_WINDOWS / "segoeui.ttf"]
    )
    for nombre in nombres:
        try:
            return ImageFont.truetype(str(nombre), tamano)
        except OSError:
            continue
    return ImageFont.load_default()


def _texto_clase(contenido_json: dict) -> str:
    partes = []
    for clave in ["titulo", "objetivo", "introduccion", "explicacion", "actividad", "resumen"]:
        partes.append(str(contenido_json.get(clave, "")))
    partes.extend(str(item) for item in contenido_json.get("ejemplos", []))
    partes.extend(str(item) for item in contenido_json.get("preguntas", []))
    return " ".join(partes).lower()


def _requiere_lamina_tablas(contenido_json: dict) -> bool:
    texto = _texto_clase(contenido_json)
    return "tabla" in texto or "multiplica" in texto


def _crear_pagina_lamina_tablas() -> Image.Image:
    imagen = Image.new("RGB", (1240, 1754), "#f4f1ea")
    dibujo = ImageDraw.Draw(imagen)
    azul = "#1f3a52"
    texto = "#232320"
    borde = "#d5d1c8"

    dibujo.rounded_rectangle((70, 70, 1170, 1684), radius=28, fill="#ffffff", outline=borde, width=3)
    dibujo.text((115, 115), "Lamina visual: tablas del 1 al 9", fill=azul, font=_fuente_pdf(54, True))
    dibujo.text((115, 185), "Recurso local generado por ProfeIA, sin costo de IA.", fill="#6c6961", font=_fuente_pdf(27))

    inicio_x, inicio_y = 115, 285
    ancho_bloque, alto_bloque = 320, 385
    for tabla in range(1, 10):
        x = inicio_x + ((tabla - 1) % 3) * 345
        y = inicio_y + ((tabla - 1) // 3) * 420
        dibujo.rounded_rectangle((x, y, x + ancho_bloque, y + alto_bloque), radius=18, fill="#f9f8f4", outline=borde, width=2)
        dibujo.text((x + 25, y + 22), f"Tabla del {tabla}", fill=azul, font=_fuente_pdf(31, True))
        for multiplicador in range(1, 10):
            fila_y = y + 78 + (multiplicador - 1) * 32
            dibujo.text(
                (x + 35, fila_y),
                f"{tabla} x {multiplicador} = {tabla * multiplicador}",
                fill=texto,
                font=_fuente_pdf(25),
            )

    dibujo.text(
        (115, 1590),
        "Actividad sugerida: elegir tres resultados, explicarlos con suma repetida y crear un problema cotidiano.",
        fill="#6c6961",
        font=_fuente_pdf(24),
    )
    return imagen


def _metadata_recurso(recurso: dict) -> dict:
    metadata = recurso.get("metadata_json") or {}
    return metadata if isinstance(metadata, dict) else {}


async def _descargar_bytes(url_o_ruta: str, cliente_http: httpx.AsyncClient) -> bytes | None:
    """
    Obtiene el contenido de un recurso ya generado, sin importar si vive en
    Supabase Storage (URL publica http/https) o en disco local (modo
    desarrollo sin Supabase configurado).
    """
    if not url_o_ruta:
        return None
    if url_o_ruta.startswith("http://") or url_o_ruta.startswith("https://"):
        try:
            respuesta = await cliente_http.get(url_o_ruta)
            respuesta.raise_for_status()
        except httpx.HTTPError:
            return None
        return respuesta.content

    ruta = Path(url_o_ruta)
    if ruta.exists() and ruta.is_file():
        return ruta.read_bytes()
    return None


def _extension_recurso(url_storage: str) -> str:
    ruta = urlparse(url_storage).path if url_storage.startswith("http") else url_storage
    return Path(ruta).suffix or ".bin"


async def _bytes_imagen_recurso(recurso: dict, cliente_http: httpx.AsyncClient) -> bytes | None:
    if recurso.get("tipo") != "imagen":
        return None
    url_storage = str(recurso.get("url_storage") or "")
    if _extension_recurso(url_storage).lower() not in EXTENSIONES_IMAGEN_VALIDAS:
        return None
    return await _descargar_bytes(url_storage, cliente_http)


def _dibujar_texto_envuelto(
    dibujo: ImageDraw.ImageDraw,
    texto: str,
    x: int,
    y: int,
    ancho_wrap: int,
    fuente,
    color: str,
    alto_linea: int = 32,
) -> int:
    for linea in wrap(str(texto), width=ancho_wrap) or [""]:
        dibujo.text((x, y), linea, fill=color, font=fuente)
        y += alto_linea
    return y


def _crear_pagina_recurso_imagen(contenido_imagen: bytes, indice: int, metadata: dict) -> Image.Image:
    pagina = Image.new("RGB", (1240, 1754), "#f4f1ea")
    dibujo = ImageDraw.Draw(pagina)
    azul = "#1f3a52"
    texto = "#232320"
    gris = "#6c6961"
    borde = "#d5d1c8"

    dibujo.rounded_rectangle((70, 70, 1170, 1684), radius=28, fill="#ffffff", outline=borde, width=3)
    dibujo.text((115, 115), f"Imagen de apoyo {indice}", fill=azul, font=_fuente_pdf(48, True))
    dibujo.text((115, 180), "Recurso visual agregado a la clase.", fill=gris, font=_fuente_pdf(25))

    with Image.open(BytesIO(contenido_imagen)) as imagen_original:
        imagen = imagen_original.convert("RGB")
        imagen.thumbnail((1010, 1040))
        x_imagen = (1240 - imagen.width) // 2
        y_imagen = 270
        dibujo.rounded_rectangle(
            (x_imagen - 14, y_imagen - 14, x_imagen + imagen.width + 14, y_imagen + imagen.height + 14),
            radius=20,
            fill="#f9f8f4",
            outline=borde,
            width=2,
        )
        pagina.paste(imagen, (x_imagen, y_imagen))

    y_texto = min(1395, y_imagen + imagen.height + 55)
    titulo = metadata.get("titulo") or f"Imagen {indice}"
    origen = metadata.get("origen") or metadata.get("fuente")
    licencia = metadata.get("licencia")

    y_texto = _dibujar_texto_envuelto(dibujo, titulo, 115, y_texto, 70, _fuente_pdf(28, True), texto, 36)
    if origen:
        y_texto += 8
        y_texto = _dibujar_texto_envuelto(
            dibujo,
            f"Fuente: {origen}",
            115,
            y_texto,
            85,
            _fuente_pdf(22),
            gris,
            29,
        )
    if licencia:
        y_texto += 4
        _dibujar_texto_envuelto(
            dibujo,
            f"Licencia: {licencia}",
            115,
            y_texto,
            85,
            _fuente_pdf(22),
            gris,
            29,
        )
    return pagina


async def _crear_paginas_recursos_imagen(
    recursos: list[dict], cliente_http: httpx.AsyncClient
) -> list[Image.Image]:
    paginas = []
    for indice, recurso in enumerate(recursos, start=1):
        contenido_imagen = await _bytes_imagen_recurso(recurso, cliente_http)
        if not contenido_imagen:
            continue
        paginas.append(
            _crear_pagina_recurso_imagen(contenido_imagen, indice, _metadata_recurso(recurso))
        )
    return paginas


def _nombre_recurso_zip(recurso: dict, indice: int) -> str:
    extension = _extension_recurso(str(recurso.get("url_storage") or ""))
    tipo = str(recurso.get("tipo") or "recurso")
    return f"recursos/{indice:02d}-{tipo}{extension}"


async def _crear_pdf_con_imagenes(
    contenido_json: dict,
    lineas: list[str],
    recursos: list[dict] | None,
    cliente_http: httpx.AsyncClient,
) -> bytes:
    paginas: list[Image.Image] = []
    ancho, alto = 1240, 1754
    margen_x, margen_y = 95, 95
    alto_linea = 34
    fuente = _fuente_pdf(24)
    fuente_titulo = _fuente_pdf(40, True)
    fuente_seccion = _fuente_pdf(29, True)

    pagina = Image.new("RGB", (ancho, alto), "#ffffff")
    dibujo = ImageDraw.Draw(pagina)
    y = margen_y

    for indice, linea in enumerate(lineas):
        texto_linea = str(linea)
        if y > alto - 120:
            paginas.append(pagina)
            pagina = Image.new("RGB", (ancho, alto), "#ffffff")
            dibujo = ImageDraw.Draw(pagina)
            y = margen_y

        if not texto_linea:
            y += alto_linea // 2
            continue

        es_titulo = indice == 0
        es_seccion = texto_linea in {
            "Objetivo",
            "Introduccion",
            "Explicacion",
            "Ejemplos",
            "Actividad",
            "Preguntas",
            "Cuestionario",
            "Tarea para el hogar",
            "Resumen",
            "Adaptacion educativa",
            "Consigna simple",
            "Rutina visual",
            "Pausas sugeridas",
            "Adaptaciones",
        }
        fuente_actual = fuente_titulo if es_titulo else fuente_seccion if es_seccion else fuente
        color = "#1f3a52" if es_titulo or es_seccion else "#232320"
        ancho_wrap = 44 if es_titulo else 78
        for sublinea in wrap(texto_linea, width=ancho_wrap) or [""]:
            if y > alto - 120:
                paginas.append(pagina)
                pagina = Image.new("RGB", (ancho, alto), "#ffffff")
                dibujo = ImageDraw.Draw(pagina)
                y = margen_y
            dibujo.text((margen_x, y), sublinea, fill=color, font=fuente_actual)
            y += 48 if es_titulo else alto_linea
        y += 14 if es_titulo or es_seccion else 8

    paginas.append(pagina)
    if _requiere_lamina_tablas(contenido_json):
        paginas.append(_crear_pagina_lamina_tablas())
    paginas.extend(await _crear_paginas_recursos_imagen(recursos or [], cliente_http))

    primera, *resto = paginas
    buffer = BytesIO()
    primera.save(buffer, "PDF", resolution=150.0, save_all=True, append_images=resto)
    return buffer.getvalue()


def _lineas_contenido(contenido_json: dict) -> list[str]:
    lineas = [
        contenido_json.get("titulo", "Clase generada"),
        "",
        "Objetivo",
        contenido_json.get("objetivo", ""),
        "",
        "Introduccion",
        contenido_json.get("introduccion", ""),
        "",
        "Explicacion",
        contenido_json.get("explicacion", ""),
        "",
        "Ejemplos",
    ]

    for ejemplo in contenido_json.get("ejemplos", []):
        lineas.append(f"- {ejemplo}")

    lineas.extend(
        [
            "",
            "Actividad",
            contenido_json.get("actividad", ""),
            "",
            "Preguntas",
        ]
    )

    for pregunta in contenido_json.get("preguntas", []):
        lineas.append(f"- {pregunta}")

    lineas.extend(["", "Cuestionario"])
    for consigna in contenido_json.get("cuestionario", []):
        lineas.append(f"- {consigna}")

    tarea_hogar = contenido_json.get("tarea_hogar")
    if tarea_hogar:
        lineas.extend(["", "Tarea para el hogar", tarea_hogar])

    lineas.extend(["", "Resumen", contenido_json.get("resumen", "")])

    apoyo = contenido_json.get("apoyo_accesibilidad") or {}
    if apoyo:
        lineas.extend(
            [
                "",
                "Adaptacion educativa",
                apoyo.get("resumen_docente", ""),
                "",
                "Consigna simple",
                apoyo.get("consigna_simple", ""),
                "",
                "Rutina visual",
            ]
        )
        lineas.extend(f"- {item}" for item in apoyo.get("rutina_visual", []))
        lineas.extend(["", "Pausas sugeridas"])
        lineas.extend(f"- {item}" for item in apoyo.get("pausas_sugeridas", []))
        lineas.extend(["", "Adaptaciones"])
        lineas.extend(f"- {item}" for item in apoyo.get("adaptaciones", []))
    return lineas


def _crear_pdf_simple(titulo: str, lineas: list[str]) -> bytes:
    """Crea un PDF basico multipagina, suficiente para la demo local."""
    lineas_pdf = []
    for linea in lineas:
        if not linea:
            lineas_pdf.append("")
            continue
        linea_texto = str(linea)
        lineas_pdf.extend(wrap(linea_texto, width=82) or [""])
        if linea_texto.startswith("- "):
            lineas_pdf.append("")

    lineas_por_pagina = 42
    paginas = [
        lineas_pdf[indice : indice + lineas_por_pagina]
        for indice in range(0, len(lineas_pdf), lineas_por_pagina)
    ] or [[]]

    objetos: list[bytes] = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
    ]

    total_paginas = len(paginas)
    pagina_objetos = [3 + indice * 2 for indice in range(total_paginas)]
    kids = " ".join(f"{numero} 0 R" for numero in pagina_objetos).encode("ascii")
    objetos.append(b"<< /Type /Pages /Kids [" + kids + b"] /Count " + str(total_paginas).encode("ascii") + b" >>")

    for indice_pagina, lineas_pagina in enumerate(paginas):
        contenido_objeto = 4 + indice_pagina * 2
        objetos.append(
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 595 842] "
            b"/Resources << /Font << /F1 "
            + str(3 + total_paginas * 2).encode("ascii")
            + b" 0 R >> >> /Contents "
            + str(contenido_objeto).encode("ascii")
            + b" 0 R >>"
        )

        comandos_texto = ["BT", "/F1 11 Tf", "50 790 Td", "14 TL"]
        for indice_linea, linea in enumerate(lineas_pagina):
            if indice_linea > 0:
                comandos_texto.append("T*")
            comandos_texto.append(f"({_limpiar_texto_pdf(linea)}) Tj")
        comandos_texto.append("ET")
        stream = "\n".join(comandos_texto).encode("latin-1", errors="replace")
        objetos.append(
            b"<< /Length "
            + str(len(stream)).encode("ascii")
            + b" >>\nstream\n"
            + stream
            + b"\nendstream"
        )

    objetos.append(
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    )

    contenido = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for numero, objeto in enumerate(objetos, start=1):
        offsets.append(len(contenido))
        contenido.extend(f"{numero} 0 obj\n".encode("ascii"))
        contenido.extend(objeto)
        contenido.extend(b"\nendobj\n")

    inicio_xref = len(contenido)
    contenido.extend(f"xref\n0 {len(objetos) + 1}\n".encode("ascii"))
    contenido.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        contenido.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    contenido.extend(
        (
            f"trailer\n<< /Size {len(objetos) + 1} /Root 1 0 R >>\n"
            f"startxref\n{inicio_xref}\n%%EOF\n"
        ).encode("ascii")
    )

    return bytes(contenido)


def _agregar_slide_texto(
    presentacion: Presentation,
    titulo: str,
    lineas: list[str],
) -> None:
    slide = presentacion.slides.add_slide(presentacion.slide_layouts[6])

    lineas_titulo = wrap(str(titulo), width=46) or [str(titulo)]
    alto_titulo = max(0.8, min(1.8, 0.42 * len(lineas_titulo) + 0.2))
    tamano_titulo = 28
    if len(lineas_titulo) == 2:
        tamano_titulo = 24
    elif len(lineas_titulo) >= 3:
        tamano_titulo = 21

    caja_titulo = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(8.9), Inches(alto_titulo))
    marco_titulo = caja_titulo.text_frame
    marco_titulo.word_wrap = True
    marco_titulo.margin_left = 0
    marco_titulo.margin_right = 0
    marco_titulo.margin_top = 0
    marco_titulo.margin_bottom = 0

    for indice, linea_titulo in enumerate(lineas_titulo):
        parrafo_titulo = (
            marco_titulo.paragraphs[0] if indice == 0 else marco_titulo.add_paragraph()
        )
        parrafo_titulo.text = linea_titulo
        parrafo_titulo.font.bold = True
        parrafo_titulo.font.size = Pt(tamano_titulo)
        parrafo_titulo.alignment = PP_ALIGN.LEFT
        parrafo_titulo.space_after = Pt(2)

    y_cuerpo = 0.35 + alto_titulo + 0.35
    alto_cuerpo = max(4.8, 7.1 - y_cuerpo)
    caja_cuerpo = slide.shapes.add_textbox(Inches(0.7), Inches(y_cuerpo), Inches(8.6), Inches(alto_cuerpo))
    marco_texto = caja_cuerpo.text_frame
    marco_texto.word_wrap = True

    for indice, linea in enumerate(lineas):
        parrafo = marco_texto.paragraphs[0] if indice == 0 else marco_texto.add_paragraph()
        parrafo.text = str(linea)
        parrafo.font.size = Pt(18)
        parrafo.space_after = Pt(8)
        if str(linea).startswith("- "):
            parrafo.level = 1


def _crear_pptx_simple(destino, contenido_json: dict) -> None:
    """Crea una presentacion editable para la demo local del MVP."""
    presentacion = Presentation()
    presentacion.slide_width = Inches(10)
    presentacion.slide_height = Inches(7.5)

    titulo = contenido_json.get("titulo", "Clase generada")
    _agregar_slide_texto(
        presentacion,
        titulo,
        [
            contenido_json.get("objetivo", ""),
            "",
            "Clase preparada con ProfeIA.",
        ],
    )
    _agregar_slide_texto(
        presentacion,
        "Introduccion",
        [contenido_json.get("introduccion", "")],
    )
    _agregar_slide_texto(
        presentacion,
        "Explicacion",
        [contenido_json.get("explicacion", "")],
    )
    _agregar_slide_texto(
        presentacion,
        "Ejemplos",
        [f"- {ejemplo}" for ejemplo in contenido_json.get("ejemplos", [])]
        or ["Sin ejemplos cargados."],
    )
    _agregar_slide_texto(
        presentacion,
        "Actividad",
        [contenido_json.get("actividad", "")],
    )
    _agregar_slide_texto(
        presentacion,
        "Preguntas",
        [f"- {pregunta}" for pregunta in contenido_json.get("preguntas", [])]
        or ["Sin preguntas cargadas."],
    )
    _agregar_slide_texto(
        presentacion,
        "Cuestionario",
        [f"- {consigna}" for consigna in contenido_json.get("cuestionario", [])]
        or ["Sin cuestionario cargado."],
    )
    _agregar_slide_texto(
        presentacion,
        "Tarea para el hogar",
        [contenido_json.get("tarea_hogar") or "Sin tarea cargada."],
    )
    _agregar_slide_texto(
        presentacion,
        "Cierre",
        [contenido_json.get("resumen", "")],
    )

    apoyo = contenido_json.get("apoyo_accesibilidad") or {}
    if apoyo:
        _agregar_slide_texto(
            presentacion,
            "Adaptacion educativa",
            [
                apoyo.get("resumen_docente", ""),
                "",
                apoyo.get("consigna_simple", ""),
            ],
        )
        _agregar_slide_texto(
            presentacion,
            "Rutina visual y pausas",
            [f"- {item}" for item in apoyo.get("rutina_visual", [])]
            + [""]
            + [f"- {item}" for item in apoyo.get("pausas_sugeridas", [])],
        )
        _agregar_slide_texto(
            presentacion,
            "Adaptaciones para el aula",
            [f"- {item}" for item in apoyo.get("adaptaciones", [])]
            or ["Sin adaptaciones generadas."],
        )

    presentacion.save(destino)


async def exportar_pdf(
    cliente,
    clase_id: UUID,
    contenido_json: dict,
    recursos: list[dict] | None = None,
) -> str:
    """
    Genera un PDF con el contenido de la clase (guion, ejemplos, actividad,
    evaluacion), lo sube a Storage y devuelve su URL publica. Antes se
    devolvia una ruta en el disco del propio backend, que solo servia para
    abrir el archivo desde la misma maquina (rompia en el celular, que no
    tiene forma de acceder al disco del servidor).
    """
    titulo = contenido_json.get("titulo", "Clase generada")
    lineas = _lineas_contenido(contenido_json)
    try:
        async with httpx.AsyncClient(timeout=30) as cliente_http:
            contenido_pdf = await _crear_pdf_con_imagenes(
                contenido_json, lineas, recursos or [], cliente_http
            )
    except Exception as error:
        print(
            "[servicio_exportacion] No se pudo crear PDF visual; "
            f"usando PDF simple. Detalle: {error}"
        )
        contenido_pdf = _crear_pdf_simple(titulo, lineas)

    resultado_storage = await guardar_recurso_docente(
        cliente=cliente,
        clase_id=clase_id,
        nombre_archivo=f"clase-{clase_id}.pdf",
        contenido=contenido_pdf,
        content_type="application/pdf",
    )
    return resultado_storage["url"]


async def exportar_pptx(cliente, clase_id: UUID, contenido_json: dict) -> str:
    """
    Genera una presentacion PPTX a partir del contenido de la clase, la
    sube a Storage y devuelve su URL publica.
    """
    buffer = BytesIO()
    _crear_pptx_simple(buffer, contenido_json)
    resultado_storage = await guardar_recurso_docente(
        cliente=cliente,
        clase_id=clase_id,
        nombre_archivo=f"clase-{clase_id}.pptx",
        contenido=buffer.getvalue(),
        content_type=(
            "application/vnd.openxmlformats-officedocument"
            ".presentationml.presentation"
        ),
    )
    return resultado_storage["url"]


async def exportar_paquete_zip(
    cliente,
    clase_id: UUID,
    contenido_json: dict,
    recursos: list[dict],
    codigo_publico: str | None,
    url_pdf: str | None,
    url_pptx: str | None,
) -> str:
    """
    Empaqueta la clase (PDF, PPTX y demas recursos) en un ZIP, lo sube a
    Storage y devuelve su URL publica. Los recursos (imagenes, PDF, PPTX)
    ya viven en Storage como URLs http(s); se descargan en memoria para
    incluirlos en el ZIP en vez de asumir que son rutas locales.
    """
    recursos_manifest = []
    for indice, recurso in enumerate(recursos, start=1):
        recursos_manifest.append(
            {
                "tipo": recurso.get("tipo"),
                "url_storage": recurso.get("url_storage"),
                "archivo_zip": _nombre_recurso_zip(recurso, indice),
                "metadata_json": recurso.get("metadata_json") or {},
            }
        )

    manifiesto = {
        "clase_id": str(clase_id),
        "codigo_publico": codigo_publico,
        "titulo": contenido_json.get("titulo", "Clase generada"),
        "resumen": contenido_json.get("resumen", ""),
        "cuestionario": contenido_json.get("cuestionario", []),
        "tarea_hogar": contenido_json.get("tarea_hogar"),
        "archivos": {
            "pdf": "clase.pdf" if url_pdf else None,
            "powerpoint": "clase.pptx" if url_pptx else None,
        },
        "recursos": recursos_manifest,
    }

    readme = (
        "Paquete ProfeIA\n\n"
        "Este ZIP contiene los materiales exportados de la clase.\n"
        "Incluye PDF, PowerPoint, manifiesto de recursos y el codigo de alumno "
        "cuando esta disponible.\n\n"
        f"Codigo alumno: {codigo_publico or 'No disponible'}\n"
    )

    buffer = BytesIO()
    async with httpx.AsyncClient(timeout=30) as cliente_http:
        contenido_pdf = await _descargar_bytes(url_pdf or "", cliente_http)
        contenido_pptx = await _descargar_bytes(url_pptx or "", cliente_http)

        with ZipFile(buffer, "w", compression=ZIP_DEFLATED) as archivo_zip:
            archivo_zip.writestr("manifest.json", json.dumps(manifiesto, ensure_ascii=False, indent=2))
            archivo_zip.writestr("README.txt", readme)

            if contenido_pdf:
                archivo_zip.writestr("clase.pdf", contenido_pdf)
            if contenido_pptx:
                archivo_zip.writestr("clase.pptx", contenido_pptx)

            for indice, recurso in enumerate(recursos, start=1):
                contenido_recurso = await _descargar_bytes(
                    str(recurso.get("url_storage") or ""), cliente_http
                )
                if not contenido_recurso:
                    continue
                archivo_zip.writestr(_nombre_recurso_zip(recurso, indice), contenido_recurso)

    resultado_storage = await guardar_recurso_docente(
        cliente=cliente,
        clase_id=clase_id,
        nombre_archivo=f"clase-{clase_id}-paquete.zip",
        contenido=buffer.getvalue(),
        content_type="application/zip",
    )
    return resultado_storage["url"]
